from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess
from typing import Dict
import random
from pathlib import Path
from nest.core import Injectable
from fastapi import UploadFile, HTTPException
from sqlalchemy.ext.asyncio import AsyncSession
import os
import httpx
from docx import Document
from datetime import datetime
from .resource_generation_model import ResourceGenerationRequest, ResourceGenerationResponse
from ..resource.resource_model import Resource
from ..resource.resource_service import ResourceService
from .image import get_image_for_presentation

@Injectable()
class ResourceGenerationService:
    def __init__(self, resource_service: ResourceService):
        self.resource_service = resource_service
        self.UPLOAD_DIR = "uploads"
        self.GENERATED_DIR = "generated_resources"
        self.PLANTILLAS_DIR = "plantillas"
        self.DOCUMENT_PROCESSING_URL = "http://pynest-app-genai:8000/document-processing/process"
        
        # Crear directorios necesarios
        for directory in [self.UPLOAD_DIR, self.GENERATED_DIR]:
            os.makedirs(directory, exist_ok=True)

    async def generate_resource(
        self,
        file: UploadFile,
        titulo: str,  # Cambiado de 'title' a 'titulo'
        modelo: str,
        diseno: int,
        tipo_recurso: str,
        cantidad_paginas: int,
        session: AsyncSession
    ) -> ResourceGenerationResponse:
        try:
            # 1. Guardar el archivo input
            input_path = os.path.join(self.UPLOAD_DIR, file.filename)
            content = await file.read()
            with open(input_path, "wb") as buffer:
                buffer.write(content)

            # 2. Procesar con document-processing
            async with httpx.AsyncClient(timeout=60.0) as client:
                with open(input_path, "rb") as f:
                    files = {"file": (file.filename, f, "application/pdf")}
                    data = {
                        "topic": titulo,  # Usando 'titulo' en lugar de 'title'
                        "model_type": modelo,
                        "temperature": 0.7
                    }
                    response = await client.post(
                        self.DOCUMENT_PROCESSING_URL,
                        files=files,
                        data=data
                    )
                    
                    if response.status_code != 200:
                        raise HTTPException(
                            status_code=response.status_code,
                            detail=f"Error en document processing: {response.text}"
                        )
                    
                    processing_result = response.json()
                    summary_file = processing_result["file_path"]

            # 3. Generar el recurso según el tipo
            if tipo_recurso == "word":
                output_path = await self._generate_word(summary_file, titulo, diseno)
            elif tipo_recurso == "ppt":
                output_path = await self._generate_ppt(summary_file, titulo, diseno, cantidad_paginas)
            else:
                output_path = await self._generate_pdf(summary_file, titulo, diseno, cantidad_paginas)

            # 4. Crear registro en resource
            resource_data = Resource(
                input_url=input_path,
                titulo=titulo,  # Usando 'titulo' en lugar de 'title'
                modelo=modelo,
                diseno=diseno,
                cantidad_paginas=cantidad_paginas,
                tipo_recurso=tipo_recurso,
                url_recurso=output_path
            )
            
            resource = await self.resource_service.create_resource(resource_data, session)

            return ResourceGenerationResponse(
                resource_id=resource.resource_id,
                summary_path=summary_file,
                resource_path=output_path
            )

        except Exception as e:
            print(f"Error en resource generation: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))
    
    async def _parse_content(self, summary_file: str) -> Dict:
        """Helper method to parse the content with markdown formatting"""
        sections = {
            'title': '',
            'slides': []
        }
        
        current_section = None
        current_content = []
        
        with open(summary_file, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                    
                # Title (##)
                if line.startswith('## '):
                    if current_section:
                        sections['slides'].append({
                            'title': current_section,
                            'content': '\n'.join(current_content)
                        })
                    sections['title'] = line[3:].strip()
                    current_section = None
                    current_content = []
                
                # Subtitle (###)
                elif line.startswith('### '):
                    if current_section:
                        sections['slides'].append({
                            'title': current_section,
                            'content': '\n'.join(current_content)
                        })
                    current_section = line[4:].strip()
                    current_content = []
                
                # Content
                elif current_section is not None:
                    # Handle bullet points
                    if line.startswith('* '):
                        line = '• ' + line[2:]
                    # Handle bold text
                    while '**' in line:
                        line = line.replace('**', '', 2)
                    current_content.append(line)
        
        # Add the last section
        if current_section:
            sections['slides'].append({
                'title': current_section,
                'content': '\n'.join(current_content)
            })
        
        return sections

    async def _generate_ppt(self, summary_file: str, titulo: str, diseno: int, cantidad_paginas: int) -> str:
        try:
            template_path = os.path.join(self.PLANTILLAS_DIR, "plantilla-ppt", f"Design-{diseno}.pptx")
            if not os.path.exists(template_path):
                raise FileNotFoundError(f"Template not found: {template_path}")

            prs = Presentation(template_path)
            sections = await self._parse_content(summary_file)
            slide_count = 0
            images = []

            # Generar dos imágenes diferentes
            print("Generando primera imagen...")
            image1 = await get_image_for_presentation(f"{titulo} - first perspective")
            if image1:
                images.append(image1)

            print("Generando segunda imagen...")
            image2 = await get_image_for_presentation(f"{titulo} - second perspective")
            if image2:
                images.append(image2)

            # Primera diapositiva (título)
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = sections['title']
            slide_count = 1

            # Diapositivas de contenido
            for i, slide_info in enumerate(sections['slides']):
                if slide_count >= cantidad_paginas:
                    break

                # Crear una diapositiva de contenido
                slide_layout = prs.slide_layouts[1] if slide_info['title'].lower() == "introducción" else prs.slide_layouts[random.choice([1, 7, 8, 9])]
                content_slide = prs.slides.add_slide(slide_layout)
                content_slide.shapes.title.text = slide_info['title']

                # Añadir contenido al slide
                body_shape = content_slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide_info['content']

                # Añadir imagen en las diapositivas 1 y 3 (segunda y cuarta página contando la portada)
                if slide_count in [1, 3] and len(images) > 0:
                    try:
                        # Usar la primera imagen para la diapositiva 1 y la segunda para la 3
                        image_index = 0 if slide_count == 1 else 1
                        if image_index < len(images):
                            content_slide.shapes.add_picture(
                                images[image_index],
                                Inches(10), Inches(5),
                                width=Inches(3), height=Inches(2)
                            )
                    except Exception as e:
                        print(f"Error adding image to slide {slide_count}: {str(e)}")

                # Formatear texto
                for paragraph in tf.paragraphs:
                    if paragraph.text.startswith('• '):
                        paragraph.level = 1
                    else:
                        paragraph.level = 0

                slide_count += 1

            # Guardar la presentación y limpiar imágenes
            output_filename = f"{titulo.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            output_path = os.path.join(self.GENERATED_DIR, output_filename)
            prs.save(output_path)

            # Limpiar imágenes generadas
            for image_path in images:
                try:
                    os.remove(image_path)
                except Exception as e:
                    print(f"Error cleaning up image {image_path}: {str(e)}")

            return output_path

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error generating PowerPoint: {str(e)}")

    async def _generate_word(self, summary_file: str, titulo: str, diseno: int) -> str:
        try:
            template_path = os.path.join(self.PLANTILLAS_DIR, "plantilla-docx", f"plantilla{diseno}.docx")
            output_filename = f"{titulo.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            output_path = os.path.join(self.GENERATED_DIR, output_filename)

            # Parse content
            sections = await self._parse_content(summary_file)

            # Create document from template
            doc = Document(template_path)

            # Replace placeholders for title and date
            for paragraph in doc.paragraphs:
                if '[titulo]' in paragraph.text:
                    paragraph.text = paragraph.text.replace('[titulo]', sections['title'])
                if '[fechaActual]' in paragraph.text:
                    paragraph.text = paragraph.text.replace('[fechaActual]', datetime.now().strftime("%d/%m/%Y"))

                # Replace introduction
                if '[Introducción]' in paragraph.text:
                    intro_section = next((s for s in sections['slides'] if 'Introducción' in s['title']), None)
                    if intro_section:
                        paragraph.text = paragraph.text.replace('[Introducción]', intro_section['content'])

                # Replace main content "Contenido Principal"
                if '[Contenido Principal]' in paragraph.text:
                    content_section = next((s for s in sections['slides'] if 'Contenido Principal' in s['title']), None)
                    if content_section:
                        paragraph.text = paragraph.text.replace('[Contenido Principal]', content_section['content'])

                # Replace applications
                if '[Aplicaciones]' in paragraph.text:
                    app_section = next((s for s in sections['slides'] if 'Aplicaciones' in s['title']), None)
                    if app_section:
                        paragraph.text = paragraph.text.replace('[Aplicaciones]', app_section['content'])

                # Replace conclusion
                if '[Conclusión]' in paragraph.text:
                    conclusion_section = next((s for s in sections['slides'] if 'Conclusión' in s['title']), None)
                    if conclusion_section:
                        paragraph.text = paragraph.text.replace('[Conclusión]', conclusion_section['content'])

            # Format text (bold for headings and bullet points)
            for paragraph in doc.paragraphs:
                if paragraph.text.startswith('**') and paragraph.text.endswith('**'):
                    paragraph.text = paragraph.text[2:-2]  # Remove ** markers
                    paragraph.runs[0].bold = True
                elif '•' in paragraph.text:
                    try:
                        paragraph.style = 'List Bullet'
                    except KeyError:
                        paragraph.style = 'Normal'  # Revert to a default style if 'List Bullet' is not available

            doc.save(output_path)
            return output_path

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error generating Word: {str(e)}")




    async def _generate_pdf(self, summary_file: str, titulo: str, diseno: int, cantidad_paginas: int) -> str:
        try:
            # First generate the PowerPoint
            pptx_path = await self._generate_ppt(summary_file, titulo, diseno, cantidad_paginas)
            
            # Create PDF directory if it doesn't exist
            pdf_dir = os.path.join(self.GENERATED_DIR, "pdf")
            os.makedirs(pdf_dir, exist_ok=True)

            # Define PDF output path
            pdf_filename = f"{titulo.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            pdf_path = os.path.join(pdf_dir, pdf_filename)

            # List of possible LibreOffice paths
            libreoffice_paths = [
                'libreoffice',
                '/usr/bin/libreoffice',
                '/usr/local/bin/libreoffice',
                'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
                'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',
            ]

            # Try converting using LibreOffice
            conversion_successful = False
            for libreoffice_path in libreoffice_paths:
                try:
                    subprocess.run(
                        [libreoffice_path, '--headless', '--convert-to', 'pdf', '--outdir', pdf_dir, pptx_path],
                        capture_output=True,
                        text=True,
                        check=True
                    )
                    conversion_successful = True
                    break
                except (subprocess.CalledProcessError, FileNotFoundError):
                    continue

            if not conversion_successful:
                raise HTTPException(
                    status_code=500,
                    detail="Failed to convert PowerPoint to PDF. Please ensure LibreOffice is installed."
                )

            # Clean up the temporary PowerPoint file if PDF was requested
            os.remove(pptx_path)

            return pdf_path

        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error generating PDF: {str(e)}")